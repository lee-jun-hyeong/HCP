#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
JSON 기반 찬양 인덱싱 시스템
"""

import json
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re

class JSONPraiseIndexer:
    """JSON 기반 찬양 인덱싱 클래스"""
    
    def __init__(self, praise_folder="Praise_PPT", output_json="praise_index.json", remove_duplicate_lines=False):
        # 리소스 경로 헬퍼: 실행파일과 같은 폴더의 파일을 찾음
        def resource_path(relative: str) -> Path:
            if getattr(sys, 'frozen', False):
                # PyInstaller로 빌드된 실행파일인 경우
                base_path = Path(sys.executable).parent
            else:
                # 개발 환경인 경우
                base_path = Path(__file__).parent
            return (base_path / relative).resolve()

        self.praise_folder = Path(praise_folder)
        self.output_json = resource_path(output_json)
        self.praise_data = []
        # 슬라이드 내 동일 라인의 중복 제거 여부 (기본: 보존)
        self.remove_duplicate_lines = remove_duplicate_lines
    
    def extract_slide_text(self, slide):
        """슬라이드에서 텍스트 추출 (슬라이드별, 줄별)"""
        slide_text = []

        def is_noise_line(text: str) -> bool:
            """가사와 무관한 잡음 텍스트(예: 'ㄴㄴㄴ', 기호 반복 등) 필터링"""
            if not text:
                return True
            # 한글 음소 자모 또는 특수문자만으로 구성된 짧은 라인 제외
            if re.fullmatch(r"[\s\-_.·•]+", text):
                return True
            # 한글 자모만으로 구성된 경우 (예: ㄴㄴㄴ, ㅁㅁ, ㅠㅠ)
            if re.fullmatch(r"[ㄱ-ㅎㅏ-ㅣ]+", text):
                return True
            # 동일 문자 3회 이상 반복만 있는 경우 (예: !!!, ---)
            if re.fullmatch(r"(.)\1{2,}", text):
                return True
            # 의미 문자가 하나도 없는 경우
            if not re.search(r"[가-힣A-Za-z0-9]", text):
                return True
            return False

        seen = set()
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue
                    if is_noise_line(text):
                        continue
                    # 필요 시에만 중복 제거
                    if self.remove_duplicate_lines:
                        key = re.sub(r"\s+", "", text.lower())
                        if key in seen:
                            continue
                        seen.add(key)
                    slide_text.append(text)

        return slide_text
    
    def extract_lyrics_from_pptx(self, file_path):
        """PPTX 파일에서 가사 추출"""
        try:
            prs = Presentation(str(file_path))
            slides_data = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = self.extract_slide_text(slide)
                if slide_text:  # 빈 슬라이드 제외
                    slides_data.append({
                        "slide_number": i + 1,
                        "text": "\n".join(slide_text),
                        "text_lines": slide_text
                    })
            
            return slides_data
            
        except Exception as e:
            print(f"[ERROR] {file_path} 처리 실패: {e}")
            return []
    
    def normalize_text(self, text):
        """검색용 텍스트 정규화"""
        if not text:
            return ""
        
        # 공백 제거, 소문자 변환
        normalized = re.sub(r'\s+', '', text.lower())
        return normalized
    
    def index_praise_files(self):
        """찬양 파일들을 JSON으로 인덱싱"""
        print("=" * 60)
        print("JSON 기반 찬양 인덱싱 시작")
        print("=" * 60)
        
        if not self.praise_folder.exists():
            print(f"[ERROR] 찬양 폴더를 찾을 수 없습니다: {self.praise_folder}")
            return False
        
        # PPTX 파일들 찾기
        pptx_files = list(self.praise_folder.glob("*.pptx"))
        print(f"발견된 PPTX 파일: {len(pptx_files)}개")
        
        for i, file_path in enumerate(pptx_files, 1):
            print(f"\n[{i}/{len(pptx_files)}] 처리 중: {file_path.name}")
            
            # 파일명에서 제목 추출
            title = file_path.stem
            
            # 가사 추출
            slides_data = self.extract_lyrics_from_pptx(file_path)
            
            if slides_data:
                # 전체 가사 텍스트 생성
                all_lyrics = []
                for slide in slides_data:
                    all_lyrics.extend(slide['text_lines'])
                full_lyrics = "\n".join(all_lyrics)
                
                # 찬양 데이터 생성
                praise_entry = {
                    "id": i,
                    "filename": file_path.name,
                    "title": title,
                    "file_path": str(file_path),
                    "lyrics": full_lyrics,
                    "slides_text": slides_data,
                    "title_normalized": self.normalize_text(title),
                    "lyrics_normalized": self.normalize_text(full_lyrics)
                }
                
                self.praise_data.append(praise_entry)
                print(f"  [OK] 성공: {len(slides_data)}개 슬라이드")
            else:
                print(f"  [FAIL] 실패: 가사 추출 불가")
        
        # JSON 파일로 저장
        self.save_to_json()
        
        print(f"\n" + "=" * 60)
        print(f"[OK] 인덱싱 완료: {len(self.praise_data)}개 찬양")
        print(f"[INFO] JSON 파일: {self.output_json}")
        print("=" * 60)
        
        return True
    
    def save_to_json(self):
        """JSON 파일로 저장"""
        try:
            with open(self.output_json, 'w', encoding='utf-8') as f:
                json.dump(self.praise_data, f, ensure_ascii=False, indent=2)
            print(f"[OK] JSON 저장 완료: {self.output_json}")
        except Exception as e:
            print(f"[ERROR] JSON 저장 실패: {e}")
    
    def load_from_json(self):
        """JSON 파일에서 로드"""
        try:
            if self.output_json.exists():
                with open(self.output_json, 'r', encoding='utf-8') as f:
                    self.praise_data = json.load(f)
                print(f"[OK] JSON 로드 완료: {len(self.praise_data)}개 찬양")
                return True
            else:
                print(f"[WARNING] JSON 파일이 없습니다: {self.output_json}")
                return False
        except Exception as e:
            print(f"[ERROR] JSON 로드 실패: {e}")
            return False
    
    def search_praises(self, query, search_type="title"):
        """찬양 검색"""
        if not self.praise_data:
            if not self.load_from_json():
                return []
        
        results = []
        query_normalized = self.normalize_text(query)
        
        for praise in self.praise_data:
            score = 0
            
            if search_type == "title":
                if query_normalized in praise['title_normalized']:
                    score = 100
                elif query in praise['title']:
                    score = 80
            elif search_type == "lyrics":
                if query_normalized in praise['lyrics_normalized']:
                    score = 100
                elif query in praise['lyrics']:
                    score = 80
            elif search_type == "both":
                if query_normalized in praise['title_normalized']:
                    score += 50
                if query_normalized in praise['lyrics_normalized']:
                    score += 50
            
            if score > 0:
                results.append({
                    'praise': praise,
                    'score': score
                })
        
        # 점수순 정렬
        results.sort(key=lambda x: x['score'], reverse=True)
        return [r['praise'] for r in results]
    
    def remove_praise_by_id(self, praise_id):
        """ID로 찬양 데이터 제거"""
        try:
            self.praise_data = [praise for praise in self.praise_data if praise['id'] != praise_id]
            print(f"[OK] 찬양 데이터 제거됨: ID {praise_id}")
            return True
        except Exception as e:
            print(f"[ERROR] 찬양 데이터 제거 실패: {e}")
            return False
    
    def add_single_file(self, file_path):
        """단일 파일 추가"""
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                print(f"[ERROR] 파일이 존재하지 않습니다: {file_path}")
                return False
            
            # 기존 데이터에서 최대 ID 찾기
            max_id = max([praise['id'] for praise in self.praise_data], default=0)
            new_id = max_id + 1
            
            # 파일명에서 제목 추출
            title = file_path.stem
            
            # 슬라이드 데이터 추출
            slides_data = self.extract_lyrics_from_pptx(file_path)
            if not slides_data:
                print(f"[WARNING] 슬라이드 데이터가 없습니다: {file_path}")
                return False
            
            # 가사 텍스트 생성
            lyrics_lines = []
            for slide in slides_data:
                lyrics_lines.extend(slide['text_lines'])
            lyrics = "\n".join(lyrics_lines)
            
            # 새 찬양 데이터 생성
            new_praise = {
                "id": new_id,
                "filename": file_path.name,
                "title": title,
                "file_path": str(file_path),
                "lyrics": lyrics,
                "slides_text": slides_data,
                "title_normalized": self.normalize_text(title),
                "lyrics_normalized": self.normalize_text(lyrics)
            }
            
            # 데이터에 추가
            self.praise_data.append(new_praise)
            print(f"[OK] 새 파일 추가됨: {title} (ID: {new_id})")
            return True
            
        except Exception as e:
            print(f"[ERROR] 파일 추가 실패: {e}")
            return False
    
    def save_to_json(self):
        """JSON 파일에 저장"""
        try:
            with open(self.output_json, 'w', encoding='utf-8') as f:
                json.dump(self.praise_data, f, ensure_ascii=False, indent=2)
            print(f"[OK] JSON 파일 저장됨: {self.output_json}")
            return True
        except Exception as e:
            print(f"[ERROR] JSON 파일 저장 실패: {e}")
            return False

def main():
    """메인 함수"""
    indexer = JSONPraiseIndexer()
    
    # 인덱싱 실행
    success = indexer.index_praise_files()
    
    if success:
        print("\n인덱싱이 완료되었습니다!")
        
        # 테스트 검색
        print("\n테스트 검색:")
        test_results = indexer.search_praises("하나님", "title")
        print(f"제목에 '하나님'이 포함된 찬양: {len(test_results)}개")
        
        for praise in test_results[:3]:  # 상위 3개만 출력
            print(f"  - {praise['title']} ({len(praise['slides_text'])}개 슬라이드)")
    else:
        print("인덱싱에 실패했습니다.")

if __name__ == "__main__":
    main()
