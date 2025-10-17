#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
완전히 수정된 JSON 기반 PPT 생성기
temp.pptx의 모든 스타일을 정확히 복제
"""

import json
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_FILL
from pptx.oxml.xmlchemy import OxmlElement
import re

class JSONPPTGeneratorFixed:
    def __init__(self, json_file="praise_index.json", template_file="temp.pptx"):
        # 리소스 경로 헬퍼: 실행파일과 같은 폴더의 파일을 찾음
        def resource_path(relative: str) -> Path:
            if getattr(sys, 'frozen', False):
                # PyInstaller로 빌드된 실행파일인 경우
                base_path = Path(sys.executable).parent
            else:
                # 개발 환경인 경우
                base_path = Path(__file__).parent
            return (base_path / relative).resolve()

        self.json_file = str(resource_path(json_file))
        self.template_file = str(resource_path(template_file))
        self.template_styles = {}
        
        # 템플릿 스타일 추출
        self.extract_template_style()
    
    def extract_template_style(self):
        """temp.pptx에서 모든 스타일 추출"""
        try:
            if not os.path.exists(self.template_file):
                print(f"[WARNING] 템플릿 파일이 없습니다: {self.template_file}")
                return
            
            template_prs = Presentation(self.template_file)
            slide = template_prs.slides[0]
            
            # 슬라이드 크기
            self.template_styles['slide_size'] = {
                'width': template_prs.slide_width,
                'height': template_prs.slide_height
            }
            
            # 배경 스타일
            self.template_styles['background'] = self.extract_background(slide)
            
            # 모든 모양 스타일
            self.template_styles['shapes'] = self.extract_all_shapes(slide)
            
            # 텍스트 스타일
            self.template_styles['text_styles'] = self.extract_text_styles(slide)
            
            print(f"[OK] 템플릿 스타일 추출 완료")
            print(f"  - 슬라이드 크기: {self.template_styles['slide_size']['width']} x {self.template_styles['slide_size']['height']}")
            print(f"  - 배경 타입: {self.template_styles['background']['type']}")
            print(f"  - 모양 수: {len(self.template_styles['shapes'])}")
            print(f"  - 텍스트 스타일 수: {len(self.template_styles['text_styles'])}")
            
        except Exception as e:
            print(f"[ERROR] 템플릿 스타일 추출 실패: {e}")
            self.template_styles = {}
    
    def extract_background(self, slide):
        """배경 스타일 추출 (이미지, 비디오, 그라디언트 지원)"""
        try:
            print(f"[DEBUG] 배경 추출 시작")
            
            # 비디오 배경 확인
            video_bg = self.extract_video_background(slide)
            if video_bg:
                return video_bg
            
            # 이미지 배경 확인
            image_bg = self.extract_image_background(slide)
            if image_bg:
                return image_bg
            
            # 일반 배경 확인
            if slide.background and slide.background.fill:
                fill = slide.background.fill
                print(f"[DEBUG] 배경 fill 타입: {getattr(fill, 'type', 'unknown')}")
                
                # 단색 배경
                if hasattr(fill, 'type') and fill.type == 1:  # SOLID
                    try:
                        if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                            color = str(fill.fore_color.rgb)
                            print(f"[DEBUG] 단색 배경 발견: {color}")
                            return {'type': 'solid', 'color': color}
                        elif hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'theme_color'):
                            theme_color = str(fill.fore_color.theme_color)
                            print(f"[DEBUG] 테마 색상 배경 발견: {theme_color}")
                            return {'type': 'theme', 'color': theme_color}
                    except Exception as e:
                        print(f"[DEBUG] 색상 추출 실패: {e}")
                
                # 그라디언트 배경
                elif hasattr(fill, 'type') and fill.type == 3:  # GRADIENT
                    print(f"[DEBUG] 그라디언트 배경 발견")
                    return {
                        'type': 'gradient',
                        'gradient_type': 'linear',
                        'angle': 0,
                        'stops': [
                            {'position': 0, 'color': '000000'},  # 검은색
                            {'position': 100, 'color': '1a1a1a'}  # 어두운 회색
                        ]
                    }
            
            print(f"[DEBUG] 기본 검은색 배경 사용")
            return {'type': 'solid', 'color': '000000'}
            
        except Exception as e:
            print(f"[WARNING] 배경 스타일 추출 실패: {e}")
            return {'type': 'solid', 'color': '000000'}
    
    def extract_video_background(self, slide):
        """비디오 배경 추출"""
        try:
            # 슬라이드의 모든 모양에서 비디오 찾기
            for shape in slide.shapes:
                if hasattr(shape, 'media_type') and shape.media_type == 2:  # 비디오
                    return {
                        'type': 'video',
                        'video_path': getattr(shape, 'media_path', None),
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height
                    }
            
            # 슬라이드 배경에서 비디오 찾기
            if hasattr(slide, 'background') and slide.background:
                try:
                    bg_xml = slide.background._element.xml
                    if 'video' in bg_xml.lower() or 'movie' in bg_xml.lower():
                        return {
                            'type': 'video',
                            'video_path': 'embedded_video',
                            'left': 0,
                            'top': 0,
                            'width': slide.slide_width,
                            'height': slide.slide_height
                        }
                except:
                    pass
            
            return None
        except Exception as e:
            print(f"[WARNING] 비디오 배경 추출 실패: {e}")
            return None
    
    def extract_image_background(self, slide):
        """이미지 배경 추출"""
        try:
            print(f"[DEBUG] 이미지 배경 추출 시작")
            
            # 슬라이드의 모든 모양에서 이미지 찾기
            for shape in slide.shapes:
                if hasattr(shape, 'image') and shape.image:
                    try:
                        image_path = getattr(shape.image, 'filename', None)
                        print(f"[DEBUG] 이미지 모양 발견: {image_path}")
                        return {
                            'type': 'image',
                            'image_path': image_path,
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height
                        }
                    except Exception as e:
                        print(f"[DEBUG] 이미지 정보 추출 실패: {e}")
            
            # 슬라이드 배경에서 이미지 찾기
            if hasattr(slide, 'background') and slide.background:
                try:
                    fill = slide.background.fill
                    if hasattr(fill, 'type') and fill.type == 2:  # PICTURE
                        print(f"[DEBUG] 배경 이미지 발견")
                        return {
                            'type': 'image',
                            'image_path': 'embedded_image',
                            'left': 0,
                            'top': 0,
                            'width': slide.slide_width,
                            'height': slide.slide_height
                        }
                except Exception as e:
                    print(f"[DEBUG] 배경 이미지 추출 실패: {e}")
            
            print(f"[DEBUG] 이미지 배경을 찾을 수 없음")
            return None
            
        except Exception as e:
            print(f"[WARNING] 이미지 배경 추출 실패: {e}")
            return None
    
    def extract_all_shapes(self, slide):
        """모든 모양 스타일 추출 (동적)"""
        shapes = []
        try:
            for shape in slide.shapes:
                shape_info = {
                    'type': type(shape).__name__,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                }
                
                # 텍스트가 있는 경우
                if hasattr(shape, 'text') and shape.text:
                    shape_info['text'] = shape.text
                    shape_info['text_frame'] = {
                        'margin_left': shape.text_frame.margin_left,
                        'margin_right': shape.text_frame.margin_right,
                        'margin_top': shape.text_frame.margin_top,
                        'margin_bottom': shape.text_frame.margin_bottom
                    }
                    
                    # 텍스트 스타일 추출
                    text_styles = []
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    style = {
                                        'text': run.text,
                                        'font_name': run.font.name,
                                        'font_size': run.font.size,
                                        'bold': run.font.bold,
                                        'italic': run.font.italic,
                                        'underline': run.font.underline,
                                        'alignment': paragraph.alignment
                                    }
                                    try:
                                        if run.font.color.rgb:
                                            style['color'] = str(run.font.color.rgb)
                                        elif run.font.color.theme_color:
                                            style['theme_color'] = run.font.color.theme_color
                                        else:
                                            style['color'] = 'FFFFFF'
                                    except:
                                        style['color'] = 'FFFFFF'
                                    text_styles.append(style)
                    shape_info['text_styles'] = text_styles
                
                # 채우기 정보
                try:
                    if hasattr(shape, 'fill') and shape.fill:
                        shape_info['fill'] = self.extract_fill_info(shape.fill)
                except:
                    pass
                
                # 선 정보
                try:
                    if hasattr(shape, 'line') and shape.line:
                        shape_info['line'] = self.extract_line_info(shape.line)
                except:
                    pass
                
                # 커넥터 타입
                if hasattr(shape, 'connector_type'):
                    shape_info['connector_type'] = shape.connector_type
                
                shapes.append(shape_info)
        except Exception as e:
            print(f"[WARNING] 모양 추출 실패: {e}")
        
        return shapes
    
    def extract_fill_info(self, fill):
        """채우기 정보 추출"""
        try:
            if hasattr(fill, 'type'):
                if fill.type == 1:  # SOLID
                    if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                        return {
                            'type': 'solid',
                            'color': str(fill.fore_color.rgb)
                        }
                elif fill.type == 2:  # PATTERN
                    return {'type': 'pattern'}
                elif fill.type == 3:  # GRADIENT
                    return {'type': 'gradient'}
            return {'type': 'none'}
        except:
            return {'type': 'none'}
    
    def extract_line_info(self, line):
        """선 정보 추출"""
        try:
            line_info = {
                'width': line.width if hasattr(line, 'width') else Pt(1)
            }
            if hasattr(line, 'color') and hasattr(line.color, 'rgb'):
                line_info['color'] = str(line.color.rgb)
            else:
                line_info['color'] = '0000FF'  # 파란색 기본값
            return line_info
        except:
            return {'width': Pt(1), 'color': '0000FF'}
    
    def extract_text_styles(self, slide):
        """텍스트 스타일 추출"""
        text_styles = []
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                style = {
                                    'text': run.text,
                                    'font_name': run.font.name,
                                    'font_size': run.font.size,
                                    'bold': run.font.bold,
                                    'italic': run.font.italic,
                                    'underline': run.font.underline,
                                    'color': 'FFFFFF',  # 흰색
                                    'alignment': paragraph.alignment
                                }
                                text_styles.append(style)
        except Exception as e:
            print(f"[WARNING] 텍스트 스타일 추출 실패: {e}")
        
        return text_styles
    
    def create_ppt_from_lyrics(self, selected_praises, output_file="merged_praises.pptx"):
        """선택된 찬양들로 PPT 생성"""
        try:
            if not self.template_styles:
                print("[ERROR] 템플릿 스타일이 없습니다")
                return False
            
            # JSON 데이터 로드
            with open(self.json_file, 'r', encoding='utf-8') as f:
                praise_data = json.load(f)
            
            # 새 프레젠테이션 생성: 템플릿을 기반으로 생성하여 테마/배경을 그대로 사용
            prs = Presentation(self.template_file)
            
            # 템플릿의 기존 슬라이드들을 모두 제거 (템플릿 내용이 포함되지 않도록)
            while len(prs.slides) > 0:
                slide_id = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(slide_id)
                del prs.slides._sldIdLst[0]
            
            print(f"[DEBUG] 템플릿 슬라이드 제거 완료, 새 슬라이드 생성 시작")
            
            # 슬라이드 크기는 템플릿에 이미 반영되어 있으므로 별도 설정 불필요
            
            # 각 찬양에 대해 슬라이드 생성
            for praise_info in selected_praises:
                praise_title = praise_info['title']
                
                # JSON에서 해당 찬양 찾기
                praise_data_item = None
                for item in praise_data:
                    if item.get('title') == praise_title:
                        praise_data_item = item
                        break
                
                if not praise_data_item:
                    print(f"[WARNING] 찬양 데이터를 찾을 수 없습니다: {praise_title}")
                    continue
                
                # 각 찬양마다 맨 앞에 빈 슬라이드 추가 (구분용)
                self.create_separator_slide(prs)
                
                # 슬라이드별로 생성
                slides_text = praise_data_item.get('slides_text', [])
                if isinstance(slides_text, str):
                    try:
                        slides_text = json.loads(slides_text)
                    except:
                        slides_text = []
                
                if not slides_text:
                    # slides_text가 없으면 전체 가사로 1개 슬라이드 생성
                    lyrics = praise_data_item.get('lyrics', '')
                    if lyrics:
                        # 제어문자/특수마커 정리 후 사용
                        lyrics_clean = self._sanitize_text(lyrics)
                        # 가사만 사용 (제목 추가하지 않음)
                        lyrics_lines = [line for line in lyrics_clean.split('\n')]
                        self.create_slide_with_style(prs, praise_title, lyrics_lines)
                else:
                    # 각 슬라이드별로 생성
                    for i, slide_text in enumerate(slides_text):
                        if isinstance(slide_text, dict) and 'text' in slide_text:
                            text_content = slide_text['text']
                        elif isinstance(slide_text, str):
                            text_content = slide_text
                        else:
                            continue
                        
                        if text_content.strip():
                            # 모든 슬라이드에 가사만 표시 (제목 추가하지 않음)
                            text_clean = self._sanitize_text(text_content)
                            self.create_slide_with_style(prs, praise_title, [line for line in text_clean.split('\n')])
            
            # PPT 저장 (재시도 로직 포함)
            import time
            max_retries = 3
            
            for attempt in range(max_retries):
                try:
                    prs.save(output_file)
                    print(f"[OK] PPT 생성 완료: {output_file}")
                    return True
                except PermissionError as e:
                    if attempt < max_retries - 1:
                        print(f"[WARNING] 파일 저장 실패 (시도 {attempt + 1}/{max_retries}): {e}")
                        time.sleep(1)  # 1초 대기 후 재시도
                    else:
                        print(f"[ERROR] 파일 저장 최종 실패: {e}")
                        # 대체 파일명으로 시도
                        import os
                        base_name = os.path.splitext(output_file)[0]
                        extension = os.path.splitext(output_file)[1]
                        timestamp = int(time.time())
                        alternative_file = f"{base_name}_{timestamp}{extension}"
                        try:
                            prs.save(alternative_file)
                            print(f"[OK] 대체 파일로 저장 완료: {alternative_file}")
                            return True
                        except Exception as alt_e:
                            print(f"[ERROR] 대체 파일 저장도 실패: {alt_e}")
                            return False
                except Exception as e:
                    print(f"[ERROR] PPT 저장 실패: {e}")
                    return False
            
        except Exception as e:
            print(f"[ERROR] PPT 생성 실패: {e}")
            import traceback
            traceback.print_exc()
            return False

    def _sanitize_text(self, text: str) -> str:
        """가사 텍스트에 섞인 특수 제어/마커를 제거·정규화한다.

        - PowerPoint 추출 시 흔한 `_x000B_` 등을 줄바꿈으로 치환
        - 혼합 개행(\r\n, \r)을 \n으로 통일
        - 수직 탭 등 제어문자를 줄바꿈으로 정규화
        - 3줄 이상 연속 개행은 2줄로 축약
        - 각 라인의 앞뒤 공백 제거
        """
        try:
            if text is None:
                return ""
            s = str(text)
            # 혼합 개행 통일
            s = s.replace("\r\n", "\n").replace("\r", "\n")
            # PPT 추출 마커 치환
            s = s.replace("_x000B_", "\n").replace("_x000C_", "\n").replace("_x000D_", "\n").replace("_x0009_", " ")
            # 제어문자 치환 (수직탭, 폼피드 등)
            s = s.replace("\u000b", "\n").replace("\u000c", "\n")
            # 다중 개행 축약
            s = re.sub(r"\n{3,}", "\n\n", s)
            # 각 라인 트리밍 (내용은 유지)
            s = "\n".join(part.strip() for part in s.split("\n"))
            return s.strip()
        except Exception:
            return str(text) if text is not None else ""
    
    def create_separator_slide(self, prs):
        """구분용 빈 슬라이드 생성"""
        try:
            # 빈 슬라이드 추가
            slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
            slide = prs.slides.add_slide(slide_layout)
            
            # 템플릿 기반 배경을 그대로 사용 (별도 적용 불필요)
            print(f"[DEBUG] 구분 슬라이드 생성 완료")
            
        except Exception as e:
            print(f"[ERROR] 구분 슬라이드 생성 실패: {e}")
    
    def create_slide_with_style(self, prs, title, lyrics_list):
        """템플릿 스타일을 적용한 슬라이드 생성"""
        try:
            # 빈 슬라이드 추가
            slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
            slide = prs.slides.add_slide(slide_layout)
            
            # 템플릿 기반 배경을 그대로 사용 (별도 적용 불필요)
            
            # 장식 요소 추가 (파란색 선들)
            self.add_decorative_elements(slide)
            
            # 가사 텍스트 박스만 추가 (제목은 가사에 포함됨)
            self.add_lyrics_textbox(slide, lyrics_list)
            
        except Exception as e:
            print(f"[ERROR] 슬라이드 생성 실패: {e}")
    
    def apply_background(self, slide):
        """배경 스타일 적용 (이미지, 비디오, 그라디언트 지원)"""
        try:
            background_info = self.template_styles['background']
            print(f"[DEBUG] 배경 적용 시작: {background_info}")
            
            bg_type = background_info.get('type', 'solid')
            
            if bg_type == 'image':
                print(f"[DEBUG] 이미지 배경 적용")
                success = self.apply_image_background(slide, background_info)
                if not success:
                    print(f"[DEBUG] 이미지 배경 적용 실패, 단색 배경으로 대체")
                    self.apply_solid_background(slide)
            elif bg_type == 'video':
                print(f"[DEBUG] 비디오 배경 적용")
                self.apply_video_background(slide)
            elif bg_type == 'gradient':
                print(f"[DEBUG] 그라데이션 배경 적용")
                self.apply_gradient_background(slide, background_info)
            elif bg_type == 'theme':
                print(f"[DEBUG] 테마 색상 배경 적용")
                self.apply_theme_background(slide, background_info)
            else:
                print(f"[DEBUG] 단색 배경 적용")
                self.apply_solid_background(slide, background_info)
                
        except Exception as e:
            print(f"[WARNING] 배경 적용 실패: {e}")
            self.apply_solid_background(slide)
    
    def apply_video_background(self, slide):
        """비디오 배경 적용"""
        try:
            video_info = self.template_styles['background']
            
            if video_info.get('video_path') and video_info['video_path'] != 'embedded_video':
                # 외부 비디오 파일이 있는 경우
                try:
                    video_shape = slide.shapes.add_movie(
                        video_info['video_path'],
                        video_info['left'], video_info['top'],
                        video_info['width'], video_info['height']
                    )
                    print(f"[OK] 비디오 배경 적용: {video_info['video_path']}")
                except Exception as e:
                    print(f"[WARNING] 비디오 파일 로드 실패: {e}")
                    # 비디오 로드 실패 시 검은색 배경으로 대체
                    self.apply_solid_background(slide)
            else:
                # 임베디드 비디오인 경우 (복사 불가)
                print(f"[WARNING] 임베디드 비디오는 복사할 수 없습니다. 단색 배경으로 대체합니다.")
                self.apply_solid_background(slide)
                
        except Exception as e:
            print(f"[WARNING] 비디오 배경 적용 실패: {e}")
            self.apply_solid_background(slide)
    
    def apply_image_background(self, slide, background_info):
        """이미지 배경 적용"""
        try:
            image_path = background_info.get('image_path')
            if not image_path or image_path == 'embedded_image':
                print(f"[DEBUG] 임베디드 이미지는 복사할 수 없음")
                return False
            
            # 외부 이미지 파일이 있는 경우
            if os.path.exists(image_path):
                try:
                    # 이미지 모양 추가
                    slide.shapes.add_picture(
                        image_path,
                        background_info.get('left', 0),
                        background_info.get('top', 0),
                        background_info.get('width', slide.slide_width),
                        background_info.get('height', slide.slide_height)
                    )
                    print(f"[OK] 이미지 배경 적용 성공: {image_path}")
                    return True
                except Exception as e:
                    print(f"[WARNING] 이미지 파일 로드 실패: {e}")
                    return False
            else:
                print(f"[WARNING] 이미지 파일을 찾을 수 없음: {image_path}")
                return False
                
        except Exception as e:
            print(f"[WARNING] 이미지 배경 적용 실패: {e}")
            return False
    
    def apply_gradient_background(self, slide, background_info):
        """그라데이션 배경 적용"""
        try:
            background = slide.background
            fill = background.fill
            fill.solid()  # 그라데이션은 복잡하므로 단색으로 대체
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검은색
            print(f"[DEBUG] 그라데이션을 단색으로 대체")
        except Exception as e:
            print(f"[WARNING] 그라데이션 배경 적용 실패: {e}")
    
    def apply_theme_background(self, slide, background_info):
        """테마 색상 배경 적용"""
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            # 테마 색상을 RGB로 변환 (기본값: 검은색)
            fill.fore_color.rgb = RGBColor(0, 0, 0)
            print(f"[DEBUG] 테마 색상을 단색으로 대체")
        except Exception as e:
            print(f"[WARNING] 테마 배경 적용 실패: {e}")
    
    def apply_solid_background(self, slide, background_info=None):
        """단색 배경 적용"""
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            
            if background_info and 'color' in background_info:
                color = background_info['color']
                try:
                    # HEX 색상을 RGB로 변환
                    if color.startswith('#'):
                        color = color[1:]
                    r = int(color[0:2], 16)
                    g = int(color[2:4], 16)
                    b = int(color[4:6], 16)
                    fill.fore_color.rgb = RGBColor(r, g, b)
                    print(f"[DEBUG] 단색 배경 적용: RGB({r}, {g}, {b})")
                except Exception as e:
                    print(f"[DEBUG] 색상 변환 실패, 검은색 사용: {e}")
                    fill.fore_color.rgb = RGBColor(0, 0, 0)
            else:
                fill.fore_color.rgb = RGBColor(0, 0, 0)  # 검은색
                print(f"[DEBUG] 기본 검은색 배경 적용")
                
        except Exception as e:
            print(f"[WARNING] 단색 배경 적용 실패: {e}")
    
    def add_decorative_elements(self, slide):
        """장식 요소 추가 (템플릿에서 동적 추출)"""
        try:
            if 'shapes' not in self.template_styles:
                return
            
            # 템플릿에서 커넥터/선 모양들 찾기
            for shape_info in self.template_styles['shapes']:
                if shape_info.get('type') == 'Connector' or shape_info.get('connector_type'):
                    try:
                        # 커넥터 재생성
                        self.recreate_shape(slide, shape_info)
                    except Exception as e:
                        print(f"[WARNING] 커넥터 재생성 실패: {e}")
        except Exception as e:
            print(f"[WARNING] 장식 요소 추가 실패: {e}")
    
    def recreate_shape(self, slide, shape_info):
        """모양 재생성"""
        try:
            if shape_info.get('type') == 'Connector' or shape_info.get('connector_type'):
                # 커넥터 재생성
                connector = slide.shapes.add_connector(
                    1,  # MSO_CONNECTOR.STRAIGHT
                    shape_info['left'], shape_info['top'],
                    shape_info['left'] + shape_info['width'], 
                    shape_info['top'] + shape_info['height']
                )
                
                # 선 스타일 적용
                if 'line' in shape_info:
                    line_info = shape_info['line']
                    if 'color' in line_info:
                        try:
                            color_hex = line_info['color'].replace('#', '')
                            r = int(color_hex[0:2], 16)
                            g = int(color_hex[2:4], 16)
                            b = int(color_hex[4:6], 16)
                            connector.line.color.rgb = RGBColor(r, g, b)
                        except:
                            connector.line.color.rgb = RGBColor(0, 0, 255)  # 파란색 기본값
                    
                    if 'width' in line_info:
                        connector.line.width = line_info['width']
                    else:
                        connector.line.width = Pt(2)
                        
        except Exception as e:
            print(f"[WARNING] 모양 재생성 실패: {e}")
    
    
    def apply_text_style(self, run, style):
        """텍스트 스타일 적용 (개선된 버전)"""
        try:
            # 폰트 이름
            if 'font_name' in style and style['font_name']:
                run.font.name = style['font_name']
                print(f"[DEBUG] 폰트 이름 적용: {style['font_name']}")
            
            # 폰트 크기
            if 'font_size' in style and style['font_size']:
                run.font.size = style['font_size']
                print(f"[DEBUG] 폰트 크기 적용: {style['font_size']}")
            
            # 폰트 스타일
            if 'bold' in style:
                run.font.bold = style['bold']
            if 'italic' in style:
                run.font.italic = style['italic']
            if 'underline' in style:
                run.font.underline = style['underline']
            
            # 색상 적용
            if 'color' in style and style['color']:
                try:
                    color = style['color']
                    if color.startswith('#'):
                        color = color[1:]
                    elif color.startswith('RGB'):
                        # RGB(255,255,255) 형태 처리
                        color = color.replace('RGB(', '').replace(')', '')
                        rgb_values = color.split(',')
                        if len(rgb_values) == 3:
                            r = int(rgb_values[0].strip())
                            g = int(rgb_values[1].strip())
                            b = int(rgb_values[2].strip())
                            run.font.color.rgb = RGBColor(r, g, b)
                            print(f"[DEBUG] RGB 색상 적용: ({r}, {g}, {b})")
                            return
                    
                    # HEX 색상 처리
                    if len(color) == 6:
                        r = int(color[0:2], 16)
                        g = int(color[2:4], 16)
                        b = int(color[4:6], 16)
                        run.font.color.rgb = RGBColor(r, g, b)
                        print(f"[DEBUG] HEX 색상 적용: RGB({r}, {g}, {b})")
                    else:
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        print(f"[DEBUG] 색상 형식 오류, 흰색 사용")
                except Exception as e:
                    print(f"[DEBUG] 색상 적용 실패: {e}, 흰색 사용")
                    run.font.color.rgb = RGBColor(255, 255, 255)
            else:
                run.font.color.rgb = RGBColor(255, 255, 255)
                print(f"[DEBUG] 기본 흰색 적용")
                
        except Exception as e:
            print(f"[WARNING] 텍스트 스타일 적용 실패: {e}")
            try:
                run.font.color.rgb = RGBColor(255, 255, 255)
            except:
                pass
    
    
    def add_lyrics_textbox(self, slide, lyrics_list):
        """가사 텍스트 박스 추가 (개선된 버전)"""
        try:
            print(f"[DEBUG] 가사 텍스트박스 추가 시작")
            
            # 템플릿에서 가사 텍스트박스 스타일 찾기
            lyrics_style = self.find_lyrics_style()
            if lyrics_style:
                print(f"[DEBUG] 템플릿 스타일 사용")
                # 템플릿 스타일 적용
                textbox = slide.shapes.add_textbox(
                    lyrics_style['left'], lyrics_style['top'],
                    lyrics_style['width'], lyrics_style['height']
                )
                text_frame = textbox.text_frame
                text_frame.clear()
                
                # 가사 텍스트 추가
                for i, lyrics in enumerate(lyrics_list):
                    if i > 0:
                        p = text_frame.add_paragraph()
                    else:
                        p = text_frame.paragraphs[0]
                    
                    run = p.runs[0] if p.runs else p.add_run()
                    run.text = lyrics
                    
                    # 템플릿 폰트 스타일 적용
                    if 'text_styles' in lyrics_style and lyrics_style['text_styles']:
                        template_style = lyrics_style['text_styles'][0]
                        self.apply_text_style(run, template_style)
                    else:
                        # 기본 스타일
                        run.font.name = '맑은 고딕'
                        run.font.size = Pt(18)
                        run.font.color.rgb = RGBColor(255, 255, 255)
                    
                    # 정렬
                    if 'text_styles' in lyrics_style and lyrics_style['text_styles']:
                        p.alignment = lyrics_style['text_styles'][0].get('alignment', PP_ALIGN.CENTER)
                    else:
                        p.alignment = PP_ALIGN.CENTER
                
                # 텍스트 프레임 마진 적용
                if 'text_frame' in lyrics_style:
                    tf = lyrics_style['text_frame']
                    text_frame.margin_left = tf.get('margin_left', 0)
                    text_frame.margin_right = tf.get('margin_right', 0)
                    text_frame.margin_top = tf.get('margin_top', 0)
                    text_frame.margin_bottom = tf.get('margin_bottom', 0)
            else:
                print(f"[DEBUG] 기본 스타일 사용")
                # 기본 가사 텍스트박스
                self.add_default_lyrics_textbox(slide, lyrics_list)
                
        except Exception as e:
            print(f"[ERROR] 가사 텍스트 박스 추가 실패: {e}")
    
    def find_lyrics_style(self):
        """템플릿에서 가사 텍스트박스 스타일 찾기"""
        try:
            if 'shapes' not in self.template_styles:
                return None
            
            # 텍스트가 있는 모양 중에서 가사용으로 적합한 것 찾기
            for shape_info in self.template_styles['shapes']:
                if 'text' in shape_info and shape_info['text'].strip():
                    # 제목이 아닌 텍스트박스 (더 큰 크기)
                    if shape_info.get('width', 0) > 5000000:  # 5cm 이상
                        print(f"[DEBUG] 가사 스타일 발견: {shape_info.get('width', 0)}")
                        return shape_info
            print(f"[DEBUG] 가사 스타일을 찾을 수 없음")
            return None
        except Exception as e:
            print(f"[WARNING] 가사 스타일 찾기 실패: {e}")
            return None
    
    def add_default_lyrics_textbox(self, slide, lyrics_list):
        """기본 가사 텍스트박스 추가 (30cm x 16cm) - 개선된 버전"""
        try:
            print(f"[DEBUG] 기본 가사 텍스트박스 생성 시작")
            
            # 슬라이드 크기 기준으로 중앙 배치
            slide_width = self.template_styles.get('slide_size', {}).get('width', 12192000)
            slide_height = self.template_styles.get('slide_size', {}).get('height', 6858000)
            
            # 30cm x 16cm 텍스트박스
            width = int(30 * 360000)  # 30cm를 EMU로 변환
            height = int(16 * 360000)  # 16cm를 EMU로 변환
            left = (slide_width - width) // 2
            top = (slide_height - height) // 2
            
            print(f"[DEBUG] 텍스트박스 크기: {width}x{height}, 위치: ({left}, {top})")
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.clear()
            
            # 가사 텍스트 추가
            for i, lyrics in enumerate(lyrics_list):
                if lyrics.strip():  # 빈 줄 제외
                    if i > 0:
                        p = text_frame.add_paragraph()
                    else:
                        p = text_frame.paragraphs[0]
                    
                    run = p.runs[0] if p.runs else p.add_run()
                    run.text = lyrics
                    
                    # 템플릿 텍스트 스타일 적용
                    template_text_styles = self.template_styles.get('text_styles', [])
                    if template_text_styles:
                        template_style = template_text_styles[0]
                        self.apply_text_style(run, template_style)
                    else:
                        # 기본 스타일
                        run.font.name = '맑은 고딕'
                        run.font.size = Pt(18)
                        run.font.color.rgb = RGBColor(255, 255, 255)
                    
                    # 정렬
                    if template_text_styles and template_text_styles[0].get('alignment'):
                        p.alignment = template_text_styles[0]['alignment']
                    else:
                        p.alignment = PP_ALIGN.CENTER
                    
                    print(f"[DEBUG] 가사 라인 추가: {lyrics[:20]}...")
            
            print(f"[DEBUG] 기본 가사 텍스트박스 생성 완료")
                
        except Exception as e:
            print(f"[ERROR] 기본 가사 텍스트 박스 추가 실패: {e}")

if __name__ == "__main__":
    # 테스트
    generator = JSONPPTGeneratorFixed(
        json_file="praise_index.json",
        template_file="temp.pptx"
    )
    
    # 테스트용 찬양 선택
    test_praises = [
        {'title': '주안에서 내 영혼'},
        {'title': '주님의 사랑'}
    ]
    
    result = generator.create_ppt_from_lyrics(test_praises, "test_fixed_style.pptx")
    if result:
        print("테스트 PPT 생성 성공!")
    else:
        print("테스트 PPT 생성 실패!")
